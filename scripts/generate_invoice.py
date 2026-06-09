import sys
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
import traceback

def replace_text_in_shape(shape, replacements):
    """Replace placeholders in a shape while preserving formatting."""
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        # Get combined text of the paragraph to check if key exists (even if split across runs)
        p_text = "".join(run.text for run in paragraph.runs)
        
        has_key = False
        for key in replacements.keys():
            if key in p_text:
                has_key = True
                break
                
        if has_key:
            # First, try to replace at individual run level to keep formatting
            for key, value in replacements.items():
                for run in paragraph.runs:
                    if key in run.text:
                        run.text = run.text.replace(key, str(value))
            
            # Re-read combined text
            p_text = "".join(run.text for run in paragraph.runs)
            
            # If any placeholder key is still present, it means it is split across runs
            for key, value in replacements.items():
                if key in p_text:
                    if paragraph.runs:
                        first_run = paragraph.runs[0]
                        first_run.text = p_text.replace(key, str(value))
                        # Clear text of all other runs in this paragraph
                        for run in paragraph.runs[1:]:
                            run.text = ""
                    p_text = "".join(run.text for run in paragraph.runs)
def get_all_shapes(shape_container):
    """Recursively collect all shapes, including nested group shapes."""
    shapes_list = []
    for shape in shape_container:
        shapes_list.append(shape)
        if shape.shape_type == 6: # Group shape
            shapes_list.extend(get_all_shapes(shape.shapes))
    return shapes_list

def convert_to_pdf(pptx_path, pdf_path):
    """Convert PPTX to PDF using PowerPoint COM interface on Windows."""
    try:
        import win32com.client
        
        # Check if PowerPoint is already running to avoid quitting the user's active app
        was_running = True
        try:
            powerpoint = win32com.client.GetActiveObject("PowerPoint.Application")
        except Exception:
            was_running = False
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        
        # Absolute paths are required for COM
        abs_pptx = str(Path(pptx_path).absolute())
        abs_pdf = str(Path(pdf_path).absolute())
        
        deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        # 32 is the constant for saving as PDF
        deck.SaveAs(abs_pdf, 32)
        deck.Close()
        
        if not was_running:
            powerpoint.Quit()
        return True
    except Exception as e:
        print(f"PDF Conversion failed: {e}")
        return False

def main():
    if len(sys.argv) < 3:
        print("Usage: python generate_invoice.py <template_pptx> <data_json> <output_pptx> [output_pdf]")
        sys.exit(1)

    template_path = sys.argv[1]
    data_path = sys.argv[2]
    output_pptx = sys.argv[3]
    output_pdf = sys.argv[4] if len(sys.argv) > 4 else None

    # Load data
    with open(data_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # Load presentation
    prs = Presentation(template_path)

    # Check which items are present in the order (up to 6 items)
    has_items = {i: (data.get(f"{{{{ITEM_{i}_NAME}}}}", "") != "") for i in range(1, 7)}

    shapes_to_delete = []

    # Process all slides and shapes
    for slide in prs.slides:
        all_shapes = get_all_shapes(slide.shapes)
        
        for shape in all_shapes:
            # Handle background strips checking
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text
                
                # Check for background strip placeholders ({{ITEM_1_BG}} to {{ITEM_6_BG}})
                # ITEM_i_BG requires item i in the order (has_items[i])
                is_bg_placeholder = False
                for i in range(1, 7):
                    bg_key = f"{{{{ITEM_{i}_BG}}}}"
                    if bg_key in text:
                        is_bg_placeholder = True
                        if not has_items.get(i, False):
                            shapes_to_delete.append(shape)
                        else:
                            shape.text_frame.text = text.replace(bg_key, "")
                        break
                
                if is_bg_placeholder:
                    continue

            # Handle standard shapes replacement
            if hasattr(shape, "text_frame") and shape.text_frame:
                replace_text_in_shape(shape, data)
            
            # Handle tables
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if hasattr(cell, "text_frame") and cell.text_frame:
                            replace_text_in_shape(cell, data)

    # Delete marked background strip shapes
    for shape in shapes_to_delete:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception as e:
            print(f"Failed to delete shape: {e}")

    # Save PPTX
    prs.save(output_pptx)
    print(f"Successfully generated PPTX: {output_pptx}")

    # Convert to PDF if requested
    if output_pdf:
        if convert_to_pdf(output_pptx, output_pdf):
            print(f"Successfully generated PDF: {output_pdf}")
        else:
            print("Failed to generate PDF. Only PPTX is available.")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"Error: {e}")
        traceback.print_exc()
        sys.exit(1)

